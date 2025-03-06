
# import streamlit as st
# import os
# from PIL import Image
# import cv2
# import pytesseract
# from docx import Document
# from pptx import Presentation
# from io import BytesIO
# import pandas as pd
# from pathlib import Path
# from langchain_community.document_loaders import PyPDFLoader, DirectoryLoader
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain_huggingface import HuggingFaceEmbeddings
# from langchain_community.vectorstores import FAISS
# from langchain.schema import Document as LangchainDocument
# from dotenv import load_dotenv, find_dotenv
# import logging

# # Configure logging
# logging.basicConfig(level=logging.INFO)

# # Load environment variables
# load_dotenv(find_dotenv())

# # Paths
# UPLOAD_FOLDER = 'data/'
# DB_FAISS_PATH = "vectorstore/db_faiss"
# os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# # Set Tesseract OCR Path
# pytesseract.pytesseract.tesseract_cmd = r"C:\\Users\\Rupesh Shinde\\Tesseract\\tesseract.exe"

# # Function to extract text from images
# def extract_text_from_image(image_path):
#     try:
#         img = cv2.imread(str(image_path))
#         if img is None:
#             logging.error(f"Failed to read image: {image_path}")
#             return ""
#         return pytesseract.image_to_string(img).strip()
#     except Exception as e:
#         logging.error(f"Error during OCR for {image_path}: {e}")
#         return ""

# # Function to load and vectorize documents
# def vectorize_documents():
#     logging.info("Starting document vectorization...")
#     documents = []

#     # Load PDFs
#     pdf_loader = DirectoryLoader(UPLOAD_FOLDER, glob="*.pdf", loader_cls=PyPDFLoader)
#     documents.extend(pdf_loader.load())

#     # Load DOCX files
#     for file in Path(UPLOAD_FOLDER).glob("*.docx"):
#         doc = Document(file)
#         text = "\n".join([para.text for para in doc.paragraphs])
#         documents.append(LangchainDocument(page_content=text, metadata={"source": file.name}))

#     # Load PPTX files
#     for file in Path(UPLOAD_FOLDER).glob("*.pptx"):
#         prs = Presentation(file)
#         for i, slide in enumerate(prs.slides):
#             text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
#             if text.strip():
#                 documents.append(LangchainDocument(page_content=text, metadata={"source": file.name, "slide": i + 1}))

#     # Load images (JPG and PNG)
#     for image_file in Path(UPLOAD_FOLDER).rglob("*.jpg"):
#         text = extract_text_from_image(image_file)
#         if text:
#             documents.append(LangchainDocument(page_content=text, metadata={"source": image_file.name}))

#     for image_file in Path(UPLOAD_FOLDER).rglob("*.png"):
#         text = extract_text_from_image(image_file)
#         if text:
#             documents.append(LangchainDocument(page_content=text, metadata={"source": image_file.name}))

#     if documents:
#         text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
#         text_chunks = text_splitter.split_documents(documents)
#         embedding_model = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
#         db = FAISS.from_documents(text_chunks, embedding_model)
#         db.save_local(DB_FAISS_PATH)
#         logging.info("Vector store updated successfully!")
#         return True  # Return success status
#     else:
#         logging.info("No new documents to vectorize.")
#         return False  # Return failure status

# # Function to handle file uploads
# def upload_files():
#     st.subheader("Upload Files")
#     uploaded_files = st.file_uploader(
#         "Choose files", type=["pdf", "docx", "pptx", "jpg", "png", "txt"], accept_multiple_files=True
#     )
#     if uploaded_files:
#         progress_bar = st.progress(0)
#         total_files = len(uploaded_files)
#         for i, uploaded_file in enumerate(uploaded_files):
#             file_path = os.path.join(UPLOAD_FOLDER, uploaded_file.name)
#             with open(file_path, "wb") as f:
#                 f.write(uploaded_file.getbuffer())
#             progress_bar.progress((i + 1) / total_files)
#             st.success(f"✅ {uploaded_file.name} uploaded successfully! ({i + 1}/{total_files})")

#         # Start vectorization
#         st.info("Vectorization started. Please wait...")
#         with st.spinner("Processing documents..."):
#             success = vectorize_documents()  # Run vectorization in the main thread
#         if success:
#             st.success("✅ Vector store updated successfully!")
#         else:
#             st.info("No new documents to vectorize.")

# # Function to display and manage uploaded files
# def manage_files():
#     st.subheader("Manage Files")
#     files = os.listdir(UPLOAD_FOLDER)
#     if files:
#         file_data = [
#             {
#                 "File Name": file,
#                 "Type": file.split(".")[-1].upper(),
#                 "Size (KB)": round(os.path.getsize(os.path.join(UPLOAD_FOLDER, file)) / 1024, 2),
#             }
#             for file in files
#         ]
#         df = pd.DataFrame(file_data)
#         st.dataframe(df, use_container_width=True)
#         selected_file = st.selectbox("Select a file to manage", files)
#         col1, col2 = st.columns(2)
#         with col1:
#             if st.button("Delete Selected File"):
#                 os.remove(os.path.join(UPLOAD_FOLDER, selected_file))
#                 st.error(f"❌ {selected_file} deleted successfully!")
#                 st.rerun()
#         with col2:
#             with open(os.path.join(UPLOAD_FOLDER, selected_file), "rb") as f:
#                 st.download_button(
#                     label="Download Selected File",
#                     data=f,
#                     file_name=selected_file,
#                     mime="application/octet-stream",
#                 )
#     else:
#         st.info("No files uploaded yet.")

# # Function to handle login
# def login():
#     st.sidebar.subheader("Login")
#     username = st.sidebar.text_input("Username")
#     password = st.sidebar.text_input("Password", type="password")
#     if st.sidebar.button("Login"):
#         with st.spinner("Logging in..."):  # Show a spinner during login
#             if username == "admin" and password == "secret":
#                 st.session_state.logged_in = True
#                 st.sidebar.success("You are logged in!")
#                 st.rerun()
#             else:
#                 st.sidebar.error("Invalid credentials.")

# # Function to handle logout
# def logout():
#     if st.sidebar.button("Logout"):
#         st.session_state.logged_in = False
#         st.sidebar.success("You have been logged out!")
#         st.rerun()

# # Main function to run the app
# def main():
#     st.title("Admin Panel")
#     if "logged_in" not in st.session_state:
#         st.session_state.logged_in = False

#     if not st.session_state.logged_in:
#         login()
#         st.stop()

#     logout()
#     upload_files()
#     manage_files()

# if __name__ == "__main__":
#     main()


#####################################################################################################################################

# import streamlit as st
# import os
# from PIL import Image
# import cv2
# import pytesseract
# from docx import Document
# from pptx import Presentation
# from io import BytesIO
# import pandas as pd
# from pathlib import Path
# from langchain_community.document_loaders import PyPDFLoader, DirectoryLoader
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain_huggingface import HuggingFaceEmbeddings
# from langchain_community.vectorstores import FAISS
# from langchain.schema import Document as LangchainDocument
# from dotenv import load_dotenv, find_dotenv
# import logging

# # Configure logging
# logging.basicConfig(level=logging.INFO)

# # Load environment variables
# load_dotenv(find_dotenv())

# # Paths
# UPLOAD_FOLDER = 'data/'
# DB_FAISS_PATH = "vectorstore/db_faiss"
# os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# # Set Tesseract OCR Path
# pytesseract.pytesseract.tesseract_cmd = r"C:\\Users\\Rupesh Shinde\\Tesseract\\tesseract.exe"

# # Function to extract text from images
# def extract_text_from_image(image_path):
#     try:
#         img = cv2.imread(str(image_path))
#         if img is None:
#             logging.error(f"Failed to read image: {image_path}")
#             return ""
#         return pytesseract.image_to_string(img).strip()
#     except Exception as e:
#         logging.error(f"Error during OCR for {image_path}: {e}")
#         return ""

# # Function to load and vectorize documents
# def vectorize_documents():
#     logging.info("Starting document vectorization...")
#     documents = []
#     # Load PDFs
#     pdf_loader = DirectoryLoader(UPLOAD_FOLDER, glob="*.pdf", loader_cls=PyPDFLoader)
#     documents.extend(pdf_loader.load())
#     # Load DOCX files
#     for file in Path(UPLOAD_FOLDER).glob("*.docx"):
#         doc = Document(file)
#         text = "\n".join([para.text for para in doc.paragraphs])
#         documents.append(LangchainDocument(page_content=text, metadata={"source": file.name}))
#     # Load PPTX files
#     for file in Path(UPLOAD_FOLDER).glob("*.pptx"):
#         prs = Presentation(file)
#         for i, slide in enumerate(prs.slides):
#             text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
#             if text.strip():
#                 documents.append(LangchainDocument(page_content=text, metadata={"source": file.name, "slide": i + 1}))
#     # Load images (JPG and PNG)
#     for image_file in Path(UPLOAD_FOLDER).rglob("*.jpg"):
#         text = extract_text_from_image(image_file)
#         if text:
#             documents.append(LangchainDocument(page_content=text, metadata={"source": image_file.name}))
#     for image_file in Path(UPLOAD_FOLDER).rglob("*.png"):
#         text = extract_text_from_image(image_file)
#         if text:
#             documents.append(LangchainDocument(page_content=text, metadata={"source": image_file.name}))
#     if documents:
#         text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
#         text_chunks = text_splitter.split_documents(documents)
#         embedding_model = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
#         db = FAISS.from_documents(text_chunks, embedding_model)
#         db.save_local(DB_FAISS_PATH)
#         logging.info("Vector store updated successfully!")
#         return True  # Return success status
#     else:
#         logging.info("No new documents to vectorize.")
#         return False  # Return failure status

# # Function to handle file uploads
# def upload_files():
#     st.subheader("Upload Files")
#     uploaded_files = st.file_uploader(
#         "Choose files", type=["pdf", "docx", "pptx", "jpg", "png", "txt"], accept_multiple_files=True
#     )
#     if uploaded_files:
#         progress_bar = st.progress(0)
#         total_files = len(uploaded_files)
#         successful_uploads = 0
#         for i, uploaded_file in enumerate(uploaded_files):
#             file_path = os.path.join(UPLOAD_FOLDER, uploaded_file.name)
#             # Check if the file already exists
#             if os.path.exists(file_path):
#                 st.warning(f"⚠️ File '{uploaded_file.name}' already exists. Skipping upload.")
#                 continue
#             # Save the file
#             with open(file_path, "wb") as f:
#                 f.write(uploaded_file.getbuffer())
#             successful_uploads += 1
#             progress_bar.progress((i + 1) / total_files)
#             st.success(f"✅ {uploaded_file.name} uploaded successfully! ({successful_uploads}/{total_files})")
#         if successful_uploads > 0:
#             # Start vectorization
#             st.info("Vectorization started. Please wait...")
#             with st.spinner("Processing documents..."):
#                 success = vectorize_documents()  # Run vectorization in the main thread
#             if success:
#                 st.success("✅ Vector store updated successfully!")
#             else:
#                 st.info("No new documents to vectorize.")
#         else:
#             st.info("No new files were uploaded.")

# # Function to display and manage uploaded files
# def manage_files():
#     st.subheader("Manage Files")
#     files = os.listdir(UPLOAD_FOLDER)
#     if files:
#         file_data = [
#             {
#                 "No.": index + 1,  # Add row number
#                 "File Name": file,
#                 "Type": file.split(".")[-1].upper(),
#                 "Size (KB)": round(os.path.getsize(os.path.join(UPLOAD_FOLDER, file)) / 1024, 2),
#             }
#             for index, file in enumerate(files)
#         ]
#         df = pd.DataFrame(file_data)
        
#         # Display files with delete buttons
#         st.write("Uploaded Files:")
#         for index, row in df.iterrows():
#             col1, col2, col3 = st.columns([1, 4, 1])  # Adjust column widths
#             with col1:
#                 st.write(row["No."])  # Display row number
#             with col2:
#                 st.write(f"{row['File Name']} ({row['Type']}, {row['Size (KB)']} KB)")  # Display file details
#             with col3:
#                 if st.button("❌", key=f"delete_{index}"):  # Cross button for deletion
#                     os.remove(os.path.join(UPLOAD_FOLDER, row["File Name"]))
#                     st.success(f"❌ {row['File Name']} deleted successfully!")
#                     st.rerun()  # Refresh the page after deletion
#     else:
#         st.info("No files uploaded yet.")

# # Function to handle login
# def login():
#     st.sidebar.subheader("Login")
#     username = st.sidebar.text_input("Username")
#     password = st.sidebar.text_input("Password", type="password")
#     if st.sidebar.button("Login"):
#         with st.spinner("Logging in..."):  # Show a spinner during login
#             if username == "admin" and password == "secret":
#                 st.session_state.logged_in = True
#                 st.sidebar.success("You are logged in!")
#                 st.rerun()
#             else:
#                 st.sidebar.error("Invalid credentials.")

# # Function to handle logout
# def logout():
#     if st.sidebar.button("Logout"):
#         st.session_state.logged_in = False
#         st.sidebar.success("You have been logged out!")
#         st.rerun()

# # Main function to run the app
# def main():
#     st.title("Admin Panel")
#     if "logged_in" not in st.session_state:
#         st.session_state.logged_in = False
#     if not st.session_state.logged_in:
#         login()
#         st.stop()
#     logout()
#     upload_files()
#     manage_files()

# if __name__ == "__main__":
#     main()


########################################################################################################################


import streamlit as st
import os
import subprocess
from PIL import Image
import cv2
import pytesseract
from docx import Document
from pptx import Presentation
import pandas as pd
from pathlib import Path
from langchain_community.document_loaders import PyPDFLoader, DirectoryLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.schema import Document as LangchainDocument
from dotenv import load_dotenv, find_dotenv
import logging
from datetime import datetime
from openpyxl import load_workbook  

# Configure logging
logging.basicConfig(level=logging.INFO)

# Load environment variables
load_dotenv(find_dotenv())

# Paths
UPLOAD_FOLDER = 'data/'
DB_FAISS_PATH = "vectorstore/db_faiss"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Set Tesseract OCR Path
pytesseract.pytesseract.tesseract_cmd = r"C:\\Users\\Rupesh Shinde\\Tesseract\\tesseract.exe"

# Functions
def extract_text_from_image(image_path):
    try:
        img = cv2.imread(str(image_path))
        if img is None:
            logging.error(f"Failed to read image: {image_path}")
            return ""
        return pytesseract.image_to_string(img).strip()
    except Exception as e:
        logging.error(f"Error during OCR for {image_path}: {e}")
        return ""

def vectorize_documents():
    logging.info("Starting document vectorization...")
    documents = []

    # Load PDFs
    pdf_loader = DirectoryLoader(UPLOAD_FOLDER, glob="*.pdf", loader_cls=PyPDFLoader)
    documents.extend(pdf_loader.load())

    # Load DOCX
    for file in Path(UPLOAD_FOLDER).glob("*.docx"):
        doc = Document(file)
        text = "\n".join([para.text for para in doc.paragraphs])
        documents.append(LangchainDocument(page_content=text, metadata={"source": file.name}))
        
    # Load Excel files
    for excel_file in Path(UPLOAD_FOLDER).glob("*.xlsx"):  # Added Excel processing
        try:
            wb = load_workbook(excel_file)
            text = ""
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    for cell in row:
                        if cell is not None:
                            text += str(cell) + " "
            documents.append(
                LangchainDocument(
                    page_content=text.strip(),
                    metadata={"source": excel_file.name}
                )
            )
        except Exception as e:
            logging.error(f"Error processing {excel_file.name}: {e}")
            continue
    
    # Load PPTX
    for file in Path(UPLOAD_FOLDER).glob("*.pptx"):
        prs = Presentation(file)
        for i, slide in enumerate(prs.slides):
            text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
            if text.strip():
                documents.append(LangchainDocument(page_content=text, metadata={"source": file.name, "slide": i + 1}))
    
    # Load Images
    image_extensions = ["*.jpg", "*.png"]
    for ext in image_extensions:
        for image_file in Path(UPLOAD_FOLDER).rglob(ext):
            text = extract_text_from_image(image_file)
            if text:
                documents.append(LangchainDocument(page_content=text, metadata={"source": image_file.name}))
    
    # Load TXT
    for txt_file in Path(UPLOAD_FOLDER).glob("*.txt"):
        try:
            with open(txt_file, 'r', encoding='utf-8') as f:
                text = f.read().strip()
            if text:
                documents.append(LangchainDocument(page_content=text, metadata={"source": txt_file.name}))
        except UnicodeDecodeError as e:
            logging.error(f"Error reading {txt_file.name}: {e}")
            continue
    
    if documents:
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
        text_chunks = text_splitter.split_documents(documents)
        embedding_model = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
        db = FAISS.from_documents(text_chunks, embedding_model)
        db.save_local(DB_FAISS_PATH)
        logging.info("Vector store updated successfully!")
        return True
    else:
        logging.info("No new documents to vectorize.")
        return False

def upload_files():
    st.subheader("Upload Files")
    uploaded_files = st.file_uploader(
        "Choose files", type=["pdf", "docx", "pptx", "jpg", "png", "txt", "xlsx"], accept_multiple_files=True
    )
    if uploaded_files:
        progress_bar = st.progress(0)
        total_files = len(uploaded_files)
        successful_uploads = 0
        for i, uf in enumerate(uploaded_files):
            file_path = os.path.join(UPLOAD_FOLDER, uf.name)
            if os.path.exists(file_path):
                st.warning(f"⚠️ File '{uf.name}' already exists. Skipping upload.")
                continue
            with open(file_path, "wb") as f:
                f.write(uf.getbuffer())
            successful_uploads += 1
            progress_bar.progress((i + 1)/total_files)
            st.success(f"✅ Uploaded: {uf.name} ({successful_uploads}/{total_files})")
        if successful_uploads:
            st.info("Vectorization started...")
            with st.spinner("Processing documents..."):
                success = vectorize_documents()
                if success:
                    st.session_state.vector_store_version = datetime.now().strftime("%Y%m%d%H%M%S")
                    st.success("Vector store updated!")
                else:
                    st.warning("No new documents to process.")

def manage_files():
    st.subheader("Manage Files")
    files = os.listdir(UPLOAD_FOLDER)
    if files:
        for idx, file_name in enumerate(files):
            file_path = os.path.join(UPLOAD_FOLDER, file_name)
            file_size = round(os.path.getsize(file_path) / 1024, 2)
            st.write(f"{idx + 1}. {file_name} (Size: {file_size} KB)")
            if st.button(f"🗑️ Delete {file_name}", key=idx):
                os.remove(file_path)
                st.success(f"{file_name} deleted!")
                st.rerun()
    else:
        st.info("No files uploaded yet.")

def login():
    st.sidebar.subheader("Login")
    username = st.sidebar.text_input("Username")
    password = st.sidebar.text_input("Password", type="password")
    if st.sidebar.button("Login"):
        with st.spinner("Logging in..."):
            if username == "admin" and password == "secret":
                st.session_state.logged_in = True
                st.sidebar.success("You are logged in!")
                st.rerun()
            else:
                st.sidebar.error("Invalid credentials.")

def logout():
    if st.sidebar.button("Logout"):
        st.session_state.logged_in = False
        st.sidebar.success("Logged out!")
        st.rerun()

def main():
    st.title("Admin Panel")
    if "logged_in" not in st.session_state:
        st.session_state.logged_in = False
    if not st.session_state.logged_in:
        login()
        st.stop()
    logout()
    upload_files()
    manage_files()

if __name__ == "__main__":
    main()
