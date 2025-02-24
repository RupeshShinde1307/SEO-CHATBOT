# import os
# from pathlib import Path
# import cv2
# import pytesseract
# from PIL import Image
# from docx import Document
# from pptx import Presentation
# from langchain_community.document_loaders import PyPDFLoader, DirectoryLoader
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# from langchain_huggingface import HuggingFaceEmbeddings
# from langchain_community.vectorstores import FAISS
# from langchain.schema import Document as LangchainDocument  # ✅ Ensure correct Document format
# from dotenv import load_dotenv, find_dotenv

# # Load environment variables
# load_dotenv(find_dotenv())

# # Paths
# DATA_PATH = "data/"
# DB_FAISS_PATH = "vectorstore/db_faiss"

# # Set Tesseract OCR Path (update this based on your installation)
# pytesseract.pytesseract.tesseract_cmd = r"C:\\Users\\Rupesh Shinde\\Tesseract\\tesseract.exe"

# # Step 1: Load Documents from Multiple Sources
# def load_documents(data_path):
#     documents = []

#     # Load PDFs
#     pdf_loader = DirectoryLoader(data_path, glob="*.pdf", loader_cls=PyPDFLoader)
#     documents.extend(pdf_loader.load())  # PDFs are already in Document format

#     # Load Word files
#     for file in Path(data_path).glob("*.docx"):
#         doc = Document(file)
#         text = "\n".join([para.text for para in doc.paragraphs])
#         documents.append(LangchainDocument(page_content=text, metadata={"source": file.name}))

#     # Load PowerPoint files
#     for file in Path(data_path).glob("*.pptx"):
#         prs = Presentation(file)
#         text = ""
#         for slide in prs.slides:
#             for shape in slide.shapes:
#                 if hasattr(shape, "text"):
#                     text += shape.text + "\n"
#         documents.append(LangchainDocument(page_content=text, metadata={"source": file.name}))

#     # Load Images (OCR)
#     for image_file in Path(data_path).glob("*.jpg"):
#         img = cv2.imread(str(image_file))
#         text = pytesseract.image_to_string(img)
#         documents.append(LangchainDocument(page_content=text, metadata={"source": image_file.name}))

#     for image_file in Path(data_path).glob("*.png"):
#         img = cv2.imread(str(image_file))
#         text = pytesseract.image_to_string(img)
#         documents.append(LangchainDocument(page_content=text, metadata={"source": image_file.name}))

#     print(f"✅ Loaded {len(documents)} documents from {data_path}")
#     return documents

# # Step 2: Create Chunks
# def create_chunks(documents):
#     text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
#     text_chunks = text_splitter.split_documents(documents)
#     print(f"✅ Created {len(text_chunks)} text chunks")
#     return text_chunks

# # Step 3: Create Vector Embeddings
# def get_embedding_model():
#     return HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")

# # Step 4: Store embeddings in FAISS
# def create_vector_store(text_chunks):
#     embedding_model = get_embedding_model()
#     print("🔄 Creating vector store...")
#     db = FAISS.from_documents(text_chunks, embedding_model)
#     db.save_local(DB_FAISS_PATH)
#     print("✅ Vector store created/updated successfully.")

# # Step 5: Main Execution
# if __name__ == "__main__":
#     print("🚀 Starting process...")
#     documents = load_documents(DATA_PATH)
#     text_chunks = create_chunks(documents)
#     create_vector_store(text_chunks)
#     print("🎉 Process completed successfully!")


import os
from pathlib import Path
import cv2
import pytesseract
from PIL import Image
from docx import Document
from pptx import Presentation
from langchain_community.document_loaders import PyPDFLoader, DirectoryLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.schema import Document as LangchainDocument
from dotenv import load_dotenv, find_dotenv

# Load environment variables
load_dotenv(find_dotenv())

# Paths
DATA_PATH = "data/"
DB_FAISS_PATH = "vectorstore/db_faiss"

# Set Tesseract OCR Path (update this based on your installation)
pytesseract.pytesseract.tesseract_cmd = r"C:\\Users\\Rupesh Shinde\\Tesseract\\tesseract.exe"

# Function to extract text from images
def extract_text_from_image(image_path):
    img = cv2.imread(str(image_path))
    if img is None:
        print(f"⚠️ Warning: Unable to read image {image_path}")
        return ""
    text = pytesseract.image_to_string(img)
    return text.strip()

# Step 1: Load Documents from Multiple Sources
def load_documents(data_path):
    documents = []

    # Load PDFs
    pdf_loader = DirectoryLoader(data_path, glob="*.pdf", loader_cls=PyPDFLoader)
    documents.extend(pdf_loader.load())

    # Load Word files
    for file in Path(data_path).glob("*.docx"):
        doc = Document(file)
        text = "\n".join([para.text for para in doc.paragraphs])
        documents.append(LangchainDocument(page_content=text, metadata={"source": file.name}))

    # Load PowerPoint files
    for file in Path(data_path).glob("*.pptx"):
        prs = Presentation(file)
        for i, slide in enumerate(prs.slides):
            text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
            if text.strip():
                documents.append(LangchainDocument(page_content=text, metadata={"source": file.name, "slide": i + 1}))

    # Load Images (OCR) - JPG and PNG
    for image_file in Path(data_path).rglob("*.jpg"):
        text = extract_text_from_image(image_file)
        if text:
            documents.append(LangchainDocument(page_content=text, metadata={"source": image_file.name}))
    
    for image_file in Path(data_path).rglob("*.png"):
        text = extract_text_from_image(image_file)
        if text:
            documents.append(LangchainDocument(page_content=text, metadata={"source": image_file.name}))

    print(f"✅ Loaded {len(documents)} documents from {data_path}")
    return documents

# Step 2: Create Chunks
def create_chunks(documents):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    text_chunks = text_splitter.split_documents(documents)
    print(f"✅ Created {len(text_chunks)} text chunks")
    return text_chunks

# Step 3: Create Vector Embeddings
def get_embedding_model():
    return HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")

# Step 4: Store embeddings in FAISS
def create_vector_store(text_chunks):
    embedding_model = get_embedding_model()
    print("🔄 Creating vector store...")
    db = FAISS.from_documents(text_chunks, embedding_model)
    db.save_local(DB_FAISS_PATH)
    print("✅ Vector store created/updated successfully.")

# Step 5: Main Execution
if __name__ == "__main__":
    print("🚀 Starting process...")
    documents = load_documents(DATA_PATH)
    text_chunks = create_chunks(documents)
    create_vector_store(text_chunks)
    print("🎉 Process completed successfully!")
